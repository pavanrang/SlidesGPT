import streamlit as st
import os
import random
import re
import requests
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from groq import Groq

def search_images(query):
    """Search for relevant images using Serper API"""
    url = "https://google.serper.dev/images"
    headers = {
        'X-API-KEY': "b98caae0101de8f452644e3e570112ddf7039737",
        'Content-Type': 'application/json'
    }
    payload = {
        'q': query,
        'num': 5
    }
    
    try:
        response = requests.post(url, headers=headers, json=payload)
        response.raise_for_status()
        results = response.json().get('images', [])
        return [img.get('imageUrl') for img in results if img.get('imageUrl')]
    except Exception as e:
        st.warning(f"Unable to fetch images: {str(e)}")
        return []

def download_image(url):
    """Download image from URL"""
    try:
        response = requests.get(url)
        response.raise_for_status()
        return io.BytesIO(response.content)
    except:
        return None

def generate_presentation_content(topic, include_images):
    """Generate detailed presentation content using Groq API"""
    client = Groq(api_key="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")
    
    image_instruction = """
    For slides that would benefit from visual aids, include an #Image_Query tag with a specific search query.
    Example:
    #Image_Query: [specific descriptive search term for relevant image]
    """ if include_images else ""
    
    prompt = f"""Create a complete, detailed presentation about {topic}. 
    The presentation should have:
    1. A title slide with subtitle
    2. An introduction/overview slide with clear objectives
    3. 4-6 detailed content slides with subtitles where appropriate
    4. A conclusion slide with key takeaways
    
    Format EXACTLY as follows (maintain these exact tags):
    #Title: [Main presentation title]
    #Subtitle: [Main presentation subtitle]
    
    #Slide: Introduction
    #Header: Introduction
    #Subheader: Overview & Objectives
    #Content: [Bullet points about what will be covered]
    {image_instruction}
    
    #Slide: [Topic 1]
    #Header: [Clear header for topic 1]
    #Subheader: [Explanatory subheader]
    #Content: [Detailed bullet points for topic 1]
    [Include sub-bullet points with - character]
    {image_instruction}
    
    [Continue with more slides]
    
    #Slide: Conclusion
    #Header: Key Takeaways
    #Subheader: Summary & Next Steps
    #Content: [Summary bullet points]
    
    Make each slide's content detailed but concise.
    Use • for main bullet points and - for sub-bullet points.
    Include relevant statistics or examples where appropriate.
    Only include #Image_Query tags for slides where visuals would significantly enhance understanding.
    """
    
    try:
        response = client.chat.completions.create(
            model="mixtral-8x7b-32768",
            messages=[
                {"role": "system", "content": "You are an expert presentation creator. Create professional, engaging, and detailed presentation content with clear structure, informative bullet points, and appropriate subtitles. Include relevant data and examples."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=4000
        )
        return response.choices[0].message.content
    except Exception as e:
        raise Exception(f"Error generating content: {str(e)}")

def parse_presentation_content(text_content):
    """Parse the generated content into structured format"""
    slides = []
    current_slide = {}
    
    lines = text_content.split('\n')
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        if not line:
            i += 1
            continue
            
        if line.startswith('#Title:'):
            current_slide = {
                'type': 'title',
                'title': line.replace('#Title:', '').strip()
            }
        elif line.startswith('#Subtitle:'):
            if current_slide.get('type') == 'title':
                current_slide['subtitle'] = line.replace('#Subtitle:', '').strip()
                slides.append(current_slide)
            
        elif line.startswith('#Slide:'):
            if current_slide and 'type' in current_slide:
                slides.append(current_slide)
            current_slide = {
                'type': 'content',
                'slide_name': line.replace('#Slide:', '').strip()
            }
            
        elif line.startswith('#Header:'):
            current_slide['header'] = line.replace('#Header:', '').strip()
            
        elif line.startswith('#Subheader:'):
            current_slide['subheader'] = line.replace('#Subheader:', '').strip()
            
        elif line.startswith('#Image_Query:'):
            current_slide['image_query'] = line.replace('#Image_Query:', '').strip()
            
        elif line.startswith('#Content:'):
            content = line.replace('#Content:', '').strip()
            # Collect all bullet points and sub-bullet points
            i += 1
            while i < len(lines) and not lines[i].strip().startswith('#'):
                if lines[i].strip():
                    content += '\n' + lines[i].strip()
                i += 1
                
            current_slide['content'] = content
            continue
            
        i += 1
    
    if current_slide and 'type' in current_slide:
        slides.append(current_slide)
    
    return slides

def add_image_to_slide(slide, image_data, position='bottom'):
    """Add image to slide at specified position"""
    if position == 'bottom':
        left = Inches(2)
        top = Inches(5)
        width = Inches(6)
        height = Inches(3)
    else:  # right
        left = Inches(6)
        top = Inches(2)
        width = Inches(4)
        height = Inches(3)
    
    try:
        slide.shapes.add_picture(image_data, left, top, width, height)
    except Exception as e:
        st.warning(f"Failed to add image to slide: {str(e)}")

def format_text_frame(text_frame, content):
    """Format the text in a shape's text frame with hierarchical bullet points"""
    if not content:
        return
        
    # Clear existing paragraphs
    for _ in range(len(text_frame.paragraphs) - 1):
        p = text_frame.paragraphs[-1]._element
        p.getparent().remove(p)
    
    first_para = text_frame.paragraphs[0]
    first_para.text = ""
    
    for point in content.split('\n'):
        point = point.strip()
        if not point:
            continue
            
        p = first_para if first_para.text == "" else text_frame.add_paragraph()
        
        # Handle sub-bullet points
        if point.startswith('-'):
            p.text = point[1:].strip()
            p.level = 1  # Set indentation level for sub-bullets
            p.font.size = Pt(16)
        else:
            p.text = point if point.startswith('•') else f'• {point}'
            p.level = 0  # Main bullet points
            p.font.size = Pt(18)
            
        p.space_after = Pt(12)
        p.space_before = Pt(6)


def add_subtitle_to_slide(slide, subtitle_text):
    """Add subtitle to a slide with proper formatting"""
    try:
        # Try to find subtitle placeholder
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text
        # Format subtitle
        text_frame = subtitle.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(20)
        paragraph.alignment = PP_ALIGN.LEFT
    except:
        # If no subtitle placeholder, create a new textbox
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(0.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.text = subtitle_text
        # Format textbox
        text_frame = textbox.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(20)
        paragraph.alignment = PP_ALIGN.LEFT

def create_ppt(slides, design_number, ppt_name, include_images):
    """Create PowerPoint presentation from structured content"""
    if not os.path.exists('GeneratedPresentations'):
        os.makedirs('GeneratedPresentations')
        
    prs = Presentation(f"Designs/Design-{design_number}.pptx")
    
    layouts = {
        'title': 0,
        'content': [1, 7, 8]
    }
    
    last_layout_index = -1
    
    for slide in slides:
        try:
            if slide['type'] == 'title':
                title_slide = prs.slides.add_slide(prs.slide_layouts[layouts['title']])
                title = title_slide.shapes.title
                title.text = slide['title']
                if 'subtitle' in slide:
                    add_subtitle_to_slide(title_slide, slide['subtitle'])
                continue
            
            layout_indices = layouts['content']
            layout_index = random.choice([i for i in layout_indices if i != last_layout_index])
            last_layout_index = layout_index
            
            content_slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
            
            title = content_slide.shapes.title
            title.text = slide.get('header', '')
            
            if 'subheader' in slide:
                add_subtitle_to_slide(content_slide, slide['subheader'])
            
            placeholder_index = 2 if layout_index == 8 else 1
            body_shape = content_slide.shapes.placeholders[placeholder_index]
            format_text_frame(body_shape.text_frame, slide.get('content', ''))
            
            # Add image if query exists and images are enabled
            if include_images and 'image_query' in slide:
                image_urls = search_images(slide['image_query'])
                if image_urls:
                    image_data = download_image(image_urls[0])
                    if image_data:
                        add_image_to_slide(content_slide, image_data)
            
        except Exception as e:
            raise Exception(f"Error creating slide '{slide.get('header', 'Unknown')}': {str(e)}")
    
    output_path = f'GeneratedPresentations/{ppt_name}.pptx'
    prs.save(output_path)
    return output_path

def main():
    st.set_page_config(page_title="AI PowerPoint Generator", layout="wide")
    
    st.title("Enhanced AI PowerPoint Generator")
    st.write("Generate professional presentations with detailed content, subtitles, and optional images")
    
    
    with st.form("presentation_form"):
        topic = st.text_area(
            "Enter the topic for the presentation:",
            help="Be specific about what you want in the presentation. Include any specific aspects you want covered."
        )
        
        design_number = st.number_input(
            "Select design template (1-7):",
            min_value=1,
            max_value=7,
            value=1,
            help="Choose a design template for your presentation"
        )
        
        include_images = st.toggle(
            "Include relevant images",
            help="When enabled, the system will search for and include relevant images in appropriate slides"
        )
        
        submitted = st.form_submit_button("Generate Presentation")
        
    if submitted and topic:
        try:
            with st.spinner("Generating presentation content..."):
                raw_content = generate_presentation_content(topic, include_images)
                
                with st.expander("View Generated Content"):
                    st.text(raw_content)
                
                slides = parse_presentation_content(raw_content)
                
                filename = re.sub(r'[^\w\s.\-\(\)]', '', topic)
                filename = filename.replace("\n", "").replace(" ", "_")
                
                with st.spinner("Creating PowerPoint presentation..."):
                    ppt_path = create_ppt(
                        slides,
                        design_number,
                        filename,
                        include_images
                    )
                
                with open(ppt_path, "rb") as f:
                    st.success("✨ Presentation generated successfully!")
                    st.download_button(
                        label="📥 Download Presentation",
                        data=f,
                        file_name=f"{filename}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                    
        except Exception as e:
            st.error(f"An error occurred: {str(e)}")
            st.error("Please try again or contact support if the problem persists.")
    
    elif submitted:
        st.warning("Please provide a topic for the presentation.")

if __name__ == "__main__":
    main()